import os
import fitz
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
import docx
from langchain_core.documents import Document


def load_code(folder_path):

    docs=[]

    for root,dirs,files in os.walk(folder_path):

        for file in files:

            file_path=os.path.join(root,file)

            ext=os.path.splitext(file_path)[1].lower()

            try:

                if ext==".pdf":

                    pdf=fitz.open(file_path)
                    has_text=False

                    for i in range(len(pdf)):

                        page=pdf.load_page(i)
                        text=page.get_text("text")

                        if text.strip():

                            has_text=True

                            docs.append(
                                Document(
                                    page_content=text,
                                    metadata={
                                        "source":file_path,
                                        "page":i+1,
                                        "type":"pdf"
                                    }
                                )
                            )

                    if not has_text:

                        images=convert_from_path(file_path)

                        for i,img in enumerate(images):

                            text=pytesseract.image_to_string(img)

                            if text.strip():

                                docs.append(
                                    Document(
                                        page_content=text,
                                        metadata={
                                            "source":file_path,
                                            "page":i+1,
                                            "type":"ocr_pdf"
                                        }
                                    )
                                )

                elif ext==".pptx":

                    ppt=Presentation(file_path)

                    for i,slide in enumerate(ppt.slides):

                        text=[]

                        for shape in slide.shapes:

                            if hasattr(shape,"text"):

                                if shape.text.strip():
                                    text.append(shape.text)

                        slide_text="\n".join(text)

                        if slide_text:

                            docs.append(
                                Document(
                                    page_content=slide_text,
                                    metadata={
                                        "source":file_path,
                                        "slide":i+1,
                                        "type":"ppt"
                                    }
                                )
                            )

                elif ext==".docx":

                    doc=docx.Document(file_path)

                    full_text=[]

                    for para in doc.paragraphs:

                        if para.text.strip():

                            full_text.append(para.text)

                    if full_text:

                        docs.append(
                            Document(
                                page_content="\n".join(full_text),
                                metadata={
                                    "source":file_path,
                                    "type":"docx"
                                }
                            )
                        )

            except Exception as e:

                print(f"Error processing {file_path}: {e}")

    print(f"Loaded {len(docs)} documents")

    return docs